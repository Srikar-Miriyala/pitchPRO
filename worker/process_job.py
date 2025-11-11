# worker/process_job.py
"""
Updated process_job with better prompt for any business idea.
"""

import os
import json
import logging
from datetime import datetime
from pathlib import Path
import re

log = logging.getLogger("process_job")
logging.basicConfig(level=logging.INFO)

from llm_client import generate as llm_generate

# Optional import for PPTX creation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

BASE_OUTPUT = Path("output")

def build_prompt_for_idea(idea: str) -> str:
    """
    Universal prompt that works for any business idea with realistic financials.
    """
    prompt = (
        f"Create a comprehensive startup pitch for: {idea}\n\n"
        
        "IMPORTANT: Return ONLY a valid JSON object with these exact keys:\n"
        "- tagline (string): Catchy one-line description\n"
        "- elevator_pitch (string): 1-2 sentence compelling description\n" 
        "- executive_summary (string): One detailed paragraph explaining the business\n"
        "- slides (array): 5-8 slides, each with 'title', 'bullets' (array of 3-5 points), 'speaker_notes'\n"
        "- financials (array): 3 years of REALISTIC projections with 'year', 'revenue', 'cost', 'profit' in INDIAN RUPEES\n"
        "- assumptions (array): 3-5 key business assumptions\n\n"
        
        "FINANCIAL REALISM GUIDELINES:\n"
        "- Year 1: Revenue ₹5L-₹50L, Costs ₹12L-₹1.2Cr, Profit negative\n"
        "- Year 2: 3-6x revenue growth from Year 1, Profit 10-30%\n"
        "- Year 3: 2-4x revenue growth from Year 2, Profit 25-50%\n"
        "- All numbers in Indian Rupees without '₹' symbol\n"
        "- Be realistic for Indian startup context\n\n"
        
        "SLIDE CONTENT GUIDELINES:\n"
        "- Slide 1: Problem/Opportunity\n"
        "- Slide 2: Our Solution\n"
        "- Slide 3: How It Works\n"
        "- Slide 4: Market Size\n"
        "- Slide 5: Business Model\n"
        "- Slide 6: Competitive Advantage\n"
        "- Slide 7: Financial Projections\n"
        "- Slide 8: Team/Ask\n\n"
        
        "FORMATTING:\n"
        "- Use **bold** for emphasis in text where appropriate\n"
        "- Make bullet points concise and impactful\n"
        "- Ensure financial projections are mathematically consistent\n\n"
        
        "Return ONLY the JSON object, no other text or explanations."
    )
    return prompt

def safe_write(path: Path, content: str, mode="w", encoding="utf-8"):
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, mode, encoding=encoding) as f:
        f.write(content)

def create_pptx_from_slides(out_dir: Path, slides: list, title="Pitch"):
    """Create PPTX from slides with robust error handling."""
    if not PPTX_AVAILABLE:
        log.info("python-pptx not installed, skipping PPTX generation")
        return None
    
    try:
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        safe_title = str(title)[:50] if title else "Business Pitch"
        title_slide.shapes.title.text = safe_title
        
        log.info("Creating PPTX with %d slides", len(slides))
        
        # Process each content slide
        valid_slides_count = 0
        for i, slide_data in enumerate(slides):
            if not isinstance(slide_data, dict):
                log.warning("Slide %d is not a dictionary, skipping", i+1)
                continue
                
            slide_title = slide_data.get('title')
            bullets = slide_data.get('bullets', [])
            
            # Skip slides without title or bullets
            if not slide_title or not bullets:
                log.warning("Slide %d missing title or bullets, skipping", i+1)
                continue
            
            # Ensure title is a string and not too long
            safe_slide_title = str(slide_title)[:100]
            
            # Use title and content layout (usually layout 1)
            if len(prs.slide_layouts) > 1:
                content_slide_layout = prs.slide_layouts[1]
            else:
                content_slide_layout = prs.slide_layouts[0]
                
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Set slide title
            try:
                slide.shapes.title.text = safe_slide_title
            except Exception as e:
                log.warning("Could not set title for slide %d: %s", i+1, e)
                continue
            
            # Add content to the slide
            try:
                # The content placeholder is usually index 1
                if len(slide.shapes.placeholders) > 1:
                    content_placeholder = slide.shapes.placeholders[1]
                    text_frame = content_placeholder.text_frame
                    
                    # Clear any existing text
                    text_frame.clear()
                    
                    # Add bullets (limit to 5 for readability)
                    for j, bullet in enumerate(bullets[:5]):
                        bullet_text = str(bullet).strip()
                        if not bullet_text:
                            continue
                            
                        if j == 0:
                            paragraph = text_frame.paragraphs[0]
                        else:
                            paragraph = text_frame.add_paragraph()
                            
                        paragraph.text = bullet_text
                        paragraph.level = 0
                        
                        # Set reasonable font size
                        for run in paragraph.runs:
                            run.font.size = Pt(16)
                
                valid_slides_count += 1
                log.debug("Added slide %d: %s", i+1, safe_slide_title)
                
            except Exception as e:
                log.warning("Failed to add content to slide %d: %s", i+1, e)
        
        # If no valid slides were created, don't save the file
        if valid_slides_count == 0:
            log.error("No valid slides could be created, skipping PPTX generation")
            return None
            
        pptx_path = out_dir / "pitch.pptx"
        prs.save(str(pptx_path))
        log.info("✅ PPTX created successfully with %d slides", valid_slides_count)
        return str(pptx_path)
        
    except Exception as e:
        log.error("❌ PPTX creation failed: %s", e)
        return None

def parse_llm_response(llm_out: str) -> dict:
    """
    Robustly parse LLM response, handling markdown, partial JSON, and various formats.
    """
    # Handle debug responses from llm_client
    if llm_out.strip().startswith('{') and any(key in llm_out for key in ['_debug', '_note']):
        try:
            debug_response = json.loads(llm_out)
            if any(key.startswith('_debug') for key in debug_response.keys()):
                return {"_debug_response": debug_response, "note": "LLM returned debug information"}
        except:
            pass

    parsed = None
    original_output = llm_out
    
    try:
        # First try direct JSON parse
        parsed = json.loads(llm_out)
        log.info("Direct JSON parse successful")
        return parsed
    except json.JSONDecodeError:
        log.info("Direct JSON parse failed, trying cleaned parsing")

    # Clean the response - remove markdown code blocks and whitespace
    cleaned_output = llm_out.strip()
    
    # Remove ```json and ``` markers
    if cleaned_output.startswith('```json'):
        cleaned_output = cleaned_output[7:].strip()
    elif cleaned_output.startswith('```'):
        cleaned_output = cleaned_output[3:].strip()
    if cleaned_output.endswith('```'):
        cleaned_output = cleaned_output[:-3].strip()
    
    # Try parsing cleaned output
    try:
        parsed = json.loads(cleaned_output)
        log.info("JSON parse successful after cleaning markdown")
        return parsed
    except json.JSONDecodeError:
        pass

    # Try to extract JSON using multiple regex patterns
    json_patterns = [
        r'\{[^{}]*(?:\{[^{}]*\}[^{}]*)*\}',  # Nested objects
        r'\{.*\}',  # Simple object
        r'\{[\s\S]*\}',  # Any content between braces
    ]
    
    for pattern in json_patterns:
        json_match = re.search(pattern, cleaned_output, re.DOTALL)
        if json_match:
            try:
                parsed = json.loads(json_match.group())
                log.info("Extracted JSON using regex pattern")
                return parsed
            except json.JSONDecodeError:
                continue

    # Try to fix common JSON issues and parse partial JSON
    try:
        # Find the start of JSON
        start_idx = original_output.find('{')
        if start_idx != -1:
            json_str = original_output[start_idx:]
            
            # Count braces to find a reasonable cutoff point
            brace_count = 0
            in_string = False
            escape_next = False
            last_valid_index = 0
            
            for i, char in enumerate(json_str):
                if escape_next:
                    escape_next = False
                    continue
                    
                if char == '\\':
                    escape_next = True
                    continue
                    
                if char == '"' and not escape_next:
                    in_string = not in_string
                    
                if not in_string:
                    if char == '{':
                        brace_count += 1
                    elif char == '}':
                        brace_count -= 1
                        if brace_count == 0:
                            last_valid_index = i
                            break
            
            if last_valid_index > 0:
                partial_json = json_str[:last_valid_index + 1]
                parsed = json.loads(partial_json)
                parsed["_note"] = "Partial response - LLM output was truncated"
                log.warning("Using partial JSON response")
                return parsed
    except Exception as e:
        log.warning("Partial JSON parsing failed: %s", e)

    # Final fallback: store structured raw response
    return {
        "raw_response": original_output, 
        "note": "Could not parse complete JSON from LLM response",
        "response_preview": original_output[:500] + "..." if len(original_output) > 500 else original_output,
        "response_length": len(original_output)
    }

def process_job(job_id: str, params: dict):
    """
    Main job processor with better error handling.
    """
    t0 = datetime.utcnow()
    out_dir = BASE_OUTPUT / job_id
    out_dir.mkdir(parents=True, exist_ok=True)

    status_path = out_dir / "status.json"
    raw_path = out_dir / "raw_llm.txt"
    pitch_path = out_dir / "pitch.json"

    # Write initial status
    safe_write(status_path, json.dumps({
        "job_id": job_id, 
        "status": "running", 
        "started_at": t0.isoformat(),
        "params": params
    }))

    idea = params.get("idea", "")
    audience = params.get("audience", "investors")
    tone = params.get("tone", "professional")
    max_tokens = params.get("max_tokens", 4000)

    prompt = build_prompt_for_idea(idea)
    log.info("Starting job %s with idea: %s", job_id, idea[:100] + "..." if len(idea) > 100 else idea)
    log.info("Prompt length: %d characters", len(prompt))

    try:
        llm_out = llm_generate(prompt, temperature=0.0, max_tokens=max_tokens)
    except Exception as e:
        error_msg = f"LLM generation failed: {e}"
        safe_write(raw_path, error_msg)
        safe_write(status_path, json.dumps({
            "job_id": job_id, 
            "status": "error", 
            "error": error_msg,
            "finished_at": datetime.utcnow().isoformat()
        }))
        log.error("LLM generation failed for job %s: %s", job_id, e)
        return {"job_id": job_id, "status": "error", "error": str(e)}

    # Save raw LLM output
    safe_write(raw_path, llm_out)
    log.info("LLM response received, length: %d characters", len(llm_out))

    # Parse the response using our robust parser
    parsed = parse_llm_response(llm_out)

    # Save pitch.json
    safe_write(pitch_path, json.dumps(parsed, indent=2, ensure_ascii=False))

    # Create PPTX if slides exist and are complete
    pptx_path = None
    if (isinstance(parsed, dict) and 
        isinstance(parsed.get('slides'), list) and 
        parsed['slides'] and
        not parsed.get('_note') and
        not parsed.get('_debug_response')):
        
        pptx_path = create_pptx_from_slides(out_dir, parsed['slides'], title=f"Pitch: {idea}")
        if pptx_path:
            log.info("PPTX created: %s", pptx_path)

    # Write final status
    finished_at = datetime.utcnow().isoformat()
    status_update = {
        "job_id": job_id, 
        "status": "done", 
        "finished_at": finished_at,
        "duration_seconds": (datetime.utcnow() - t0).total_seconds(),
        "has_pptx": bool(pptx_path)
    }
    
    safe_write(status_path, json.dumps(status_update, indent=2))
    
    # Log completion summary
    if isinstance(parsed, dict) and parsed.get('tagline'):
        log.info("Job %s completed successfully with tagline: %s", job_id, parsed['tagline'])
    else:
        log.info("Job %s completed with results", job_id)
    
    return {
        "job_id": job_id, 
        "out_dir": str(out_dir),
        "status": "done",
        "has_pptx": bool(pptx_path)
    }

if __name__ == "__main__":
    # Test the job processor
    test_params = {
        "idea": "AI-powered education platform for rural students",
        "audience": "investors", 
        "tone": "professional",
        "max_tokens": 4000
    }
    result = process_job("test-local", test_params)
    print("Test result:", result)