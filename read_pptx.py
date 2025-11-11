# read_pptx.py
from pptx import Presentation
import sys

path = "output/2920268216d6/pitch.pptx"

try:
    prs = Presentation(path)
    print("Successfully opened:", path)
    print("Slides:", len(prs.slides))
    for i, s in enumerate(prs.slides, start=1):
        texts = []
        for sh in s.shapes:
            try:
                if hasattr(sh, "text") and sh.text and sh.text.strip():
                    texts.append(sh.text.strip().replace("\n", " / "))
            except Exception:
                continue
        print(f"Slide {i}: first texts ->", texts[:3])
except Exception as e:
    print("Error reading PPTX:", repr(e))
    sys.exit(2)
