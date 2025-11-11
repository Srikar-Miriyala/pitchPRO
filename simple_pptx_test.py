# simple_pptx_test.py
from pptx import Presentation
from pptx.util import Inches

def create_simple_pptx():
    """Create a very simple PPTX to test if python-pptx works."""
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Test Presentation"
    
    # Content slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Test Slide"
    slide.shapes.placeholders[1].text = "This is a test bullet point"
    
    # Save
    prs.save("test_simple.pptx")
    print("✅ Simple PPTX created: test_simple.pptx")

if __name__ == "__main__":
    create_simple_pptx()