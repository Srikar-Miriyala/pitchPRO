# utils/pptx_builder.py
# Safer PPTX builder: uses blank layouts + explicit textboxes + image for financials.
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pathlib import Path
from typing import Dict, Any, List
import matplotlib.pyplot as plt
import tempfile
import os

def sanitize(s: str) -> str:
    if s is None:
        return ""
    if not isinstance(s, str):
        s = str(s)
    # Replace common problematic unicode sequences
    return (
        s.replace("\u2014", "-")
         .replace("\u2013", "-")
         .replace("\u2018", "'")
         .replace("\u2019", "'")
         .replace("\u201c", '"')
         .replace("\u201d", '"')
         .replace("â", "")  # artifact seen in your JSON
    )

def _add_textbox(slide, left, top, width, height, text, font_size=18, bold=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = sanitize(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    return tb

def _add_bullets(slide, left, top, width, height, bullets: List[str], bullet_size=18):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = sanitize(b)
        p.level = 0
        p.font.size = Pt(bullet_size)
    return tb

def _render_financials_to_image(fin_rows: List[Dict[str, Any]], out_image_path: str):
    # Create a simple table image using matplotlib
    # fin_rows: list of dicts with keys year, revenue, cost, profit
    if not fin_rows:
        # create a small placeholder image
        fig, ax = plt.subplots(figsize=(6, 2))
        ax.text(0.5, 0.5, "No financials provided", ha='center', va='center', fontsize=14)
        ax.axis('off')
        fig.tight_layout()
        fig.savefig(out_image_path, dpi=150, bbox_inches='tight')
        plt.close(fig)
        return out_image_path

    headers = ["Year", "Revenue (INR Lakhs)", "Cost (INR Lakhs)", "Profit (INR Lakhs)"]
    table_data = []
    for r in fin_rows:
        table_data.append([r.get("year", ""), r.get("revenue", ""), r.get("cost", ""), r.get("profit", "")])

    fig, ax = plt.subplots(figsize=(6, 2 + 0.4 * len(table_data)))
    ax.axis('off')
    table = ax.table(cellText=table_data, colLabels=headers, loc='center', cellLoc='center')
    table.auto_set_font_size(False)
    table.set_fontsize(10)
    table.scale(1, 1.2)
    fig.tight_layout()
    fig.savefig(out_image_path, dpi=150, bbox_inches='tight')
    plt.close(fig)
    return out_image_path

def build_pptx(pitch_obj: Dict[str, Any], out_path: Path):
    """
    pitch_obj keys: tagline, elevator_pitch, executive_summary, slides[], financials[], assumptions[]
    This builder writes a safer PPTX (blank layouts, textboxes, image for financials).
    """
    prs = Presentation()
    # Title slide: keep simple
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(1.0), pitch_obj.get("tagline", ""), font_size=30, bold=True)
    _add_textbox(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(1.0), pitch_obj.get("elevator_pitch", "")[:300], font_size=18)

    slides = pitch_obj.get("slides", [])
    for idx, s in enumerate(slides[:7], start=1):
        title = s.get("title", f"Slide {idx}")
        bullets = s.get("bullets", ["TBD"])
        notes = s.get("speaker_notes", "")

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        # Title textbox
        _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.8), title, font_size=26, bold=True)
        # Bullets area
        _add_bullets(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(4.5), bullets, bullet_size=18)
        # Speaker notes as small textbox at bottom (safer than notes_slide)
        if notes:
            _add_textbox(slide, Inches(0.5), Inches(5.6), Inches(9), Inches(0.8), "Notes: " + notes[:400], font_size=12)

        # If this is the financials slide, add an image of the table below bullets
        if "financial" in title.lower() or "financials" in title.lower() or idx == 6:
            fin_rows = pitch_obj.get("financials", [])
            # create temp image
            tmpdir = tempfile.gettempdir()
            img_path = os.path.join(tmpdir, f"pitch_fin_{os.getpid()}_{idx}.png")
            _render_financials_to_image(fin_rows, img_path)
            # add picture centered below bullets
            try:
                slide.shapes.add_picture(img_path, Inches(1.0), Inches(2.5), width=Inches(8.0))
            except Exception:
                # fallback: ignore image if insertion fails
                pass
            finally:
                # try removing temp file if exists
                try:
                    if os.path.exists(img_path):
                        os.remove(img_path)
                except Exception:
                    pass

    # Add an appendix slide for assumptions & exec summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.8), "Appendix", font_size=26, bold=True)
    exec_sum = pitch_obj.get("executive_summary", "")
    assumptions = pitch_obj.get("assumptions", [])
    combined = ("Executive Summary:\n" + exec_sum + "\n\nAssumptions:\n" + "\n".join(assumptions))[:2500]
    _add_textbox(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(5.0), combined, font_size=12)

    # Save output
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
