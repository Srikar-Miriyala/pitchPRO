# tools/build_safe_pptx.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

def safe_textbox(slide, left, top, width, height, text, font_size=20):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    return txBox

def sanitize(s: str) -> str:
    if not isinstance(s, str):
        return str(s)
    return s.replace("\u2014", "-").replace("\u2013", "-").replace("\u2018","'").replace("\u2019","'").replace("â", "")

def build_safe_pptx(out_path):
    prs = Presentation()
    # build 6 safe slides (blank layout + explicit textboxes)
    for i in range(1, 7):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = f"Safe Slide {i}"
        bullets = [f"Point {j} for slide {i}" for j in range(1,4)]
        safe_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1), sanitize(title), font_size=28)
        safe_textbox(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(4), "\n".join(bullets), font_size=18)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print("Wrote safe pptx:", out_path)

if __name__ == "__main__":
    import pathlib
    out = pathlib.Path("output/safe_test/pitch_safe.pptx")
    build_safe_pptx(out)
